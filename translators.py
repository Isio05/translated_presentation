# Import builtin libs
import os
from os import path
from glob import glob
from shutil import copyfile
from threading import Thread
from queue import Queue
from abc import ABC, abstractmethod
# Import third-party libs
import zipfile
from bs4 import BeautifulSoup as bs
from pptx import Presentation
import boto3
# Import custom libs
from shared_variables import AWS_ACCESS_KEY_ID, AWS_SECRET_ACCESS_KEY
from utils import *


class Translator(ABC):
    input_l = "pl"
    output_l = "en"

    @abstractmethod
    def process_specific_file(self):
        pass

    def __init__(self, file_to_translate: str = None, prepare_target_files: bool = True):
        # Naming stuff
        self.old_extension = None
        self.file_to_translate = file_to_translate
        self.file_to_translate_zipped = ""
        # Archive files
        self.archive_source = None
        self.archive_target = None
        # Methods required on start-up
        self.translate_service = None
        self.connect_to_translate_service()
        if prepare_target_files:
            self.convert_file_ext()
            self.open_zips()

    @classmethod
    def change_input_language(cls, new_input_l):
        cls.input_l = new_input_l

    @classmethod
    def change_ouput_language(cls, new_output_l):
        cls.output_l = new_output_l

    def connect_to_translate_service(self):
        self.translate_service = boto3.client(service_name='translate', region_name='us-east-1', use_ssl=True,
                                              aws_access_key_id=AWS_ACCESS_KEY_ID,
                                              aws_secret_access_key=AWS_SECRET_ACCESS_KEY
                                              )

    def open_zips(self):
        # Source archive should be opened just in read mode, no modifications are applied on i
        self.archive_source = zipfile.ZipFile(path.join(SOURCE_FOLDER, self.file_to_translate_zipped), "r")

        # Target archive is starting as complete copy of source archive, opened in write mode,
        # it's located in the target folder and thus can have the same name
        copyfile(path.join(SOURCE_FOLDER, self.file_to_translate_zipped),
                 path.join(TARGET_FOLDER, self.file_to_translate_zipped))
        self.archive_target = zipfile.ZipFile(path.join(TARGET_FOLDER, self.file_to_translate_zipped), "w")

    def request_translation(self, text_input):
        # Requests are made with usage of translate client prepared with class initialization
        result = self.translate_service.translate_text(Text=text_input,
                                                       SourceLanguageCode=self.input_l,
                                                       TargetLanguageCode=self.output_l)
        return result['TranslatedText']

    def convert_file_ext(self, mode="change_source"):
        """
        Changes the extension of file to zip and backwards
        :param mode: str["change_source", "final_change"]
        """
        # Rename using original absolute path and that path with modified extension
        if mode == "change_source":
            # Old extension is returned for usage in backwards conversion
            self.old_extension = self.file_to_translate.rsplit(".")[1]
            self.file_to_translate_zipped = self.file_to_translate.rsplit(".")[0] + ".zip"
            os.rename(path.join(SOURCE_FOLDER, self.file_to_translate), path.join(SOURCE_FOLDER, self.file_to_translate_zipped))
        # Rename using original absolute path and that path with modified extension
        elif mode == "final_change":
            os.rename(path.join(SOURCE_FOLDER, self.file_to_translate_zipped), path.join(SOURCE_FOLDER, self.file_to_translate))
            os.rename(path.join(TARGET_FOLDER, self.file_to_translate_zipped), path.join(TARGET_FOLDER, self.file_to_translate))
        else:
            raise RuntimeError("Unknown file extension change mode.")

    def threaded_text_translation(self, texts_to_translate):
        # Each worker is in fact a subscript for the list of text to translate
        # Each thread iterates through given text in the list, sends it to API and saves translation to dictionary
        translation = {}
        q = Queue()

        def threader():
            while True:
                worker = q.get()
                translated = self.request_translation(text_input=texts_to_translate[worker])
                translation[texts_to_translate[worker]] = translated
                q.task_done()

        # Ten threads are spawned
        for i in range(10):
            t = Thread(target=threader)
            t.daemon = True
            t.start()

        # Each possibile index for the list of text to translate is put in the queue
        for text_pos in range(len(texts_to_translate)):
            q.put(text_pos)

        q.join()

        return translation


class PresentationTranslator(Translator):
    def __init__(self, file_to_translate: str = None):
        super().__init__(file_to_translate, False)
        self.num_of_slides = 0
        self.user_num_of_slides = None

    @staticmethod
    def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
        if len(paragraph.runs) == 0:
            return None
        p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
        # remove all but the first run
        for idx, run in enumerate(paragraph.runs):
            if idx == 0:
                continue
            p.remove(run._r)
        paragraph.runs[0].text = new_text

    def process_specific_file(self):
        prs = Presentation(path.join(SOURCE_FOLDER, self.file_to_translate))
        already_translated = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text_frame = shape.text_frame
                for par in text_frame.paragraphs:
                    if len(text_frame.text) == 0 or text_frame.text in already_translated:
                        continue
                    translated_text = self.request_translation(text_frame.text)
                    self.replace_paragraph_text_retaining_initial_formatting(par, translated_text)
                    already_translated.append(text_frame.text)

        prs.save(path.join(TARGET_FOLDER, self.file_to_translate))


class DocumentTranslator(Translator):
    def __init__(self, file_to_translate: str = None):
        super().__init__(file_to_translate, True)

    def process_specific_file(self):
        """Opens file contained in zip file without extraction"""
        # Relative paths used to navigate within xml file
        contents_file_location = "word/"
        contents_file = "document.xml"
        contents_file_rel_path = contents_file_location + contents_file

        # Rewrite each file separately from the source archive to containing translation
        for item in self.archive_source.infolist():
            buffer = self.archive_source.read(item.filename)
            if not item.filename.startswith("word/document"):
                self.archive_target.writestr(item, buffer)

        # Open it directly from archive and bs starts to look for text to translate
        # After the translation the text is overwritten with translation
        # The end of the loop is to write into archive_2
        current_slide_data = self.archive_source.read(contents_file_rel_path)
        xml_soup = bs(current_slide_data.decode("UTF-8"), 'lxml')

        # Text is surrounded by "w:t"
        # After finding it is converted to strings
        texts_to_translate = xml_soup.find_all("w:t")
        texts_to_translate = [text.string for text in texts_to_translate]
        translation = self.threaded_text_translation(texts_to_translate)

        # The source slide is unpacked into simple string
        # Using dictionary that contains translations and sources, text will be replaced in the string
        # After the operation, string is encoded to basic format
        current_slide_data_decoded = current_slide_data.decode("UTF-8")
        for item, definition in translation.items():
            current_slide_data_decoded = current_slide_data_decoded.replace("<w:t>" + item,
                                                                            "<w:t>" + definition)
            current_slide_data_decoded = current_slide_data_decoded.replace('<w:t xml:space="preserve">' + item,
                                                                            '<w:t xml:space="preserve">' + definition)
        current_slide_data_encoded = current_slide_data_decoded.encode("UTF-8")

        # To write into archive, the source file must exist
        # The "temp" folder will contain ready to write xmls with translations
        # After the operation the folder temp is removed
        # Using created temp folder, create there xml file containing translation
        # Subsequently, write (wb) that file to translation archive
        f = open(path.join(TEMP_FOLDER, contents_file), "wb")
        f.write(current_slide_data_encoded)
        f.close()
        self.archive_target.write(path.join(TEMP_FOLDER, contents_file), contents_file_rel_path)

        self.archive_source.close()
        self.archive_target.close()

        self.convert_file_ext("final_change")


class WorkbookTranslator(Translator):
    def __init__(self, file_to_translate: str = None):
        super().__init__(file_to_translate, False)

    def process_specific_file(self):
        return 0


def menu():
    # Left for testing purposes
    file = "lite.pptx" # input("Type in file type with extension or 'exit': ")

    while True:
        if file == "exit":
            break
        file_type = path.splitext(file)[1]
        if file_type == ".docx":
            translate = DocumentTranslator(file_to_translate=file)
            translate.process_specific_file()
        elif file_type == ".pptx":
            translate = PresentationTranslator(file_to_translate=file)
            translate.process_specific_file()
        elif file_type == ".xlsx":
            translate = WorkbookTranslator(file_to_translate=file)
            translate.process_specific_file()
        else:
            print("Wrong file extension")

        file = "exit"


def translate_folder():
    folder = input("Set folder located in the script folder: ")

    for extension in ALLOWED_EXTENSIONS:
        # Glob requires absolute path to list files of given extension
        files = glob(folder + "\\**\*.{}".format(extension), recursive=True)
        # Program is prepared to work with folders/files located in the same directory as the script
        files_rels = [folder.split("\\\\")[-1] + x.replace(folder, "") for x in files]
        for file in files_rels:
            file_type = path.splitext(file)[1]
            if file_type == ".docx":
                translate = DocumentTranslator(file_to_translate=file)
                translate.process_specific_file()
            elif file_type == ".pptx":
                translate = PresentationTranslator(file_to_translate=file)
                translate.process_specific_file()
            elif file_type == ".xlsx":
                translate = WorkbookTranslator(file_to_translate=file)
                translate.process_specific_file()
            os.remove(path.join(path.dirname(__file__), file))

menu()
